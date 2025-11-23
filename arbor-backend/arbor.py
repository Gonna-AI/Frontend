import streamlit as st
import os
import sys
import json
import warnings
from datetime import datetime, timedelta
from pathlib import Path
from typing import List, Dict, Any, Optional
import pandas as pd
import re
import math
from collections import defaultdict

# Import from setup.py
try:
    from setup import (
        ArborConfig,
        GoogleServicesAuth,
        GoogleDriveConnector,
        GoogleCalendarConnector,
        GoogleSheetsConnector,
        EmailAutomation,
        VectorStoreManager,
        ComprehensiveTools,
        AdvancedTools,
        EnhancedAutonomousAgent,
        arbor_do,
        ask_arbor,
        get_todays_schedule,
        read_tasks,
        add_task,
        email_to,
        system_status,
        refresh_all_data,
        llm,
        embeddings,
        vectorstore,
        qa_chain,
        comprehensive_tools,
        advanced_tools,
        enhanced_arbor_agent,
        google_auth,
        gdrive_connector,
        calendar_connector,
        sheets_connector,
        email_automation,
        vectorstore_manager,
        all_documents
    )
    ARBOR_AVAILABLE = True
except ImportError as e:
    st.error(f"Could not import from setup.py: {e}")
    ARBOR_AVAILABLE = False

# Additional imports for Claims Intelligence
try:
    import nltk
    from nltk.corpus import stopwords
    from nltk.tokenize import word_tokenize, sent_tokenize
    from nltk.stem import PorterStemmer
    import fitz  # PyMuPDF
    from docx import Document as DocxDocument
    from pptx import Presentation
    from sentence_transformers import SentenceTransformer
    from sklearn.metrics.pairwise import cosine_similarity
    import numpy as np
    from transformers import pipeline
    
    # Download NLTK data if needed
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)
    try:
        nltk.data.find('corpora/stopwords')
    except LookupError:
        nltk.download('stopwords', quiet=True)
    try:
        nltk.data.find('corpora/wordnet')
    except LookupError:
        nltk.download('wordnet', quiet=True)
    
    SEARCH_AVAILABLE = True
except ImportError as e:
    SEARCH_AVAILABLE = False
    IMPORT_ERROR = str(e)

warnings.filterwarnings('ignore')

# Page config
st.set_page_config(
    page_title="Arbor AI Assistant",
    page_icon="🌳",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS with ClerkTree branding
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
    
    * {
        font-family: 'Inter', sans-serif;
    }
    
    .stApp {
        background: linear-gradient(135deg, #0a0a1f 0%, #1a1a3f 50%, #0f0f2e 100%);
        background-attachment: fixed;
    }
    
    .stApp::before {
        content: '';
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        background-image: 
            radial-gradient(2px 2px at 20% 30%, rgba(79, 70, 229, 0.1), transparent),
            radial-gradient(2px 2px at 60% 70%, rgba(139, 92, 246, 0.1), transparent),
            radial-gradient(1px 1px at 50% 50%, rgba(59, 130, 246, 0.08), transparent);
        background-size: 200% 200%;
        animation: drift 20s ease infinite;
        pointer-events: none;
    }
    
    @keyframes drift {
        0%, 100% { background-position: 0% 0%; }
        50% { background-position: 100% 100%; }
    }
    
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1a3f 0%, #0f0f2e 100%);
        border-right: 1px solid rgba(79, 70, 229, 0.3);
    }
    
    .main-header {
        background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
        padding: 2.5rem;
        border-radius: 20px;
        margin-bottom: 2rem;
        box-shadow: 0 20px 60px rgba(79, 70, 229, 0.4);
        text-align: center;
        position: relative;
        overflow: hidden;
    }
    
    .main-header::before {
        content: '';
        position: absolute;
        top: -50%;
        left: -50%;
        width: 200%;
        height: 200%;
        background: radial-gradient(circle, rgba(255,255,255,0.1) 0%, transparent 70%);
        animation: rotate 10s linear infinite;
    }
    
    @keyframes rotate {
        0% { transform: rotate(0deg); }
        100% { transform: rotate(360deg); }
    }
    
    .main-header h1 {
        color: white;
        font-size: 3rem;
        font-weight: 700;
        margin: 0;
        position: relative;
        z-index: 1;
        text-shadow: 0 0 30px rgba(255,255,255,0.5);
    }
    
    .main-header p {
        color: rgba(255,255,255,0.95);
        font-size: 1.2rem;
        margin-top: 0.5rem;
        position: relative;
        z-index: 1;
    }
    
    .stChatMessage {
        background: rgba(26, 26, 63, 0.7);
        backdrop-filter: blur(15px);
        border: 1px solid rgba(79, 70, 229, 0.3);
        border-radius: 16px;
        padding: 20px;
        margin: 15px 0;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.4);
        transition: all 0.3s ease;
    }
    
    .stChatMessage:hover {
        transform: translateY(-2px);
        box-shadow: 0 12px 40px rgba(79, 70, 229, 0.3);
        border-color: rgba(79, 70, 229, 0.5);
    }
    
[data-testid="stChatInput"] {
        background: transparent !important;
        border: none !important;
        padding: 0 !important;
        position: sticky !important;
        bottom: 0 !important;
        z-index: 100 !important;
        display: flex !important;
        align-items: center !important;
        min-height: 70px !important;
    }
    
    [data-testid="stChatInput"] > div {
        background: rgba(26, 26, 63, 0.8) !important;
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
        border-radius: 12px !important;
        padding: 0 !important;
        display: flex !important;
        align-items: center !important;
        min-height: 56px !important;
        width: 100% !important;
    }
    
    [data-testid="stChatInput"]:focus-within > div {
        border-color: rgba(79, 70, 229, 0.8) !important;
        box-shadow: 0 0 20px rgba(79, 70, 229, 0.4) !important;
        background: rgba(30, 30, 40, 1) !important;
    }
    
    [data-testid="stChatInput"] textarea,
    .stChatInput textarea,
    [data-testid="stChatInput"] textarea:active,
    [data-testid="stChatInput"] textarea:focus,
    [data-testid="stChatInput"] textarea:hover,
    [data-testid="stChatInput"] textarea[aria-invalid="true"],
    [data-testid="stChatInput"] textarea[aria-invalid="false"],
    [data-testid="stChatInput"] textarea:invalid,
    [data-testid="stChatInput"] textarea:user-invalid {
        background: transparent !important;
        color: #e0e7ff !important;
        border: none !important;
        border-radius: 12px !important;
        padding: 14px 20px !important;
        font-size: 15px !important;
        transition: all 0.3s ease !important;
        outline: none !important;
        box-shadow: none !important;
        resize: none !important;
        flex: 1 !important;
        min-height: 52px !important;
        max-height: 52px !important;
        line-height: 1.5 !important;
        padding-top: 16px !important;
    }
    
    [data-testid="stChatInput"] textarea::placeholder {
        color: #94a3b8 !important;
    }
    
    [data-testid="stChatInput"] textarea:focus::placeholder {
        color: #cbd5e1 !important;
    }
    
    [data-testid="stChatInput"] button,
    [data-testid="stChatInputSubmitButton"],
    [data-testid="stChatInput"] [data-testid="baseButton-header"],
    [data-testid="stChatInput"] [data-testid="baseButton-headerNoPadding"] {
        background: transparent !important;
        border: none !important;
        padding: 0 !important;
        margin: 0 12px 0 0 !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        transition: all 0.3s ease !important;
        border-radius: 8px !important;
        min-height: 40px !important;
        height: 40px !important;
        width: 40px !important;
        flex-shrink: 0 !important;
        position: relative !important;
        align-self: center !important;
    }
    
    [data-testid="stChatInput"] > div > div:last-child {
        display: flex !important;
        align-items: center !important;
        height: 100% !important;
    }
    
    [data-testid="stChatInput"] button:hover,
    [data-testid="stChatInputSubmitButton"]:hover {
        background: rgba(79, 70, 229, 0.1) !important;
    }
    
    [data-testid="stChatInput"] button svg,
    [data-testid="stChatInputSubmitButton"] svg {
        color: rgba(79, 70, 229, 0.8) !important;
        transition: all 0.3s ease !important;
    }
    
    [data-testid="stChatInput"] button:hover svg,
    [data-testid="stChatInputSubmitButton"]:hover svg {
        color: rgba(79, 70, 229, 1) !important;
    }

    .stTextInput > div > div > input,
    .stTextInput input,
    input[type="text"],
    div[data-baseweb="input"] input {
        border-radius: 10px !important;
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
        background: rgba(26, 26, 63, 0.8) !important;
        color: #e0e7ff !important;
        padding: 1rem 1.5rem !important;
        font-size: 1rem !important;
        transition: all 0.3s ease !important;
    }
    
    .stTextInput > div > div > input:focus,
    .stTextInput input:focus,
    input[type="text"]:focus,
    div[data-baseweb="input"] input:focus {
        border-color: rgba(79, 70, 229, 0.8) !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.2) !important;
        outline: none !important;
    }
    
    .stTextInput > div > div > input:invalid,
    .stTextInput > div > div > input:user-invalid,
    .stTextInput > div > div > input[aria-invalid="true"],
    .stTextInput input:invalid,
    .stTextInput input:user-invalid,
    input[type="text"]:invalid,
    input[type="text"]:user-invalid,
    div[data-baseweb="input"] input:invalid,
    div[data-baseweb="input"] input:user-invalid {
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
    }
    
    .stTextInput > div > div > input:focus:invalid,
    .stTextInput > div > div > input:focus:user-invalid,
    .stTextInput > div > div > input:focus[aria-invalid="true"],
    .stTextInput input:focus:invalid,
    input[type="text"]:focus:invalid {
        border: 2px solid rgba(79, 70, 229, 0.8) !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.2) !important;
    }
    
    .stTextInput > div[data-baseweb="input"],
    div[data-baseweb="input"] {
        border: none !important;
        border-radius: 10px !important;
    }
    
    .stTextInput > div[data-baseweb="input"] > div,
    div[data-baseweb="input"] > div {
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
        border-radius: 10px !important;
    }
    
    .stTextInput > div[data-baseweb="input"]:focus-within > div,
    div[data-baseweb="input"]:focus-within > div {
        border: 2px solid rgba(79, 70, 229, 0.8) !important;
        border-radius: 10px !important;
    }
    
    .stSelectbox > div > div,
    div[data-baseweb="select"] > div {
        border-radius: 10px !important;
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
        background: rgba(26, 26, 63, 0.8) !important;
        transition: all 0.3s ease !important;
    }
    
    .stSelectbox > div > div:focus-within,
    div[data-baseweb="select"]:focus-within > div {
        border-color: rgba(79, 70, 229, 0.8) !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.2) !important;
    }
    
    .stSelectbox [data-baseweb="select"] > div[aria-invalid="true"],
    .stSelectbox [data-baseweb="select"] > div:invalid {
        border: 2px solid rgba(79, 70, 229, 0.4) !important;
    }
    
    .stSelectbox [data-baseweb="select"]:focus-within > div[aria-invalid="true"],
    .stSelectbox [data-baseweb="select"]:focus-within > div:invalid {
        border: 2px solid rgba(79, 70, 229, 0.8) !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.2) !important;
    }
    
    .stButton button {
        background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
        color: white;
        border-radius: 12px;
        border: none;
        padding: 12px 28px;
        font-weight: 600;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(79, 70, 229, 0.4);
    }
    
    .stButton button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 25px rgba(79, 70, 229, 0.6);
    }
    
    .stat-card {
        background: linear-gradient(135deg, rgba(79, 70, 229, 0.15) 0%, rgba(124, 58, 237, 0.15) 100%);
        padding: 1.5rem;
        border-radius: 15px;
        border: 1px solid rgba(79, 70, 229, 0.3);
        text-align: center;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.2);
    }
    
    .stat-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 8px 30px rgba(79, 70, 229, 0.3);
        border-color: rgba(79, 70, 229, 0.6);
    }
    
    .stat-number {
        font-size: 2.5rem;
        font-weight: 700;
        background: linear-gradient(135deg, #60a5fa 0%, #a78bfa 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
    }
    
    .stat-label {
        color: #cbd5e1;
        font-size: 0.9rem;
        margin-top: 0.5rem;
        font-weight: 500;
    }
    
    .mode-button {
        background: rgba(79, 70, 229, 0.2);
        border: 2px solid rgba(79, 70, 229, 0.4);
        border-radius: 12px;
        padding: 15px;
        margin: 10px 0;
        cursor: pointer;
        transition: all 0.3s ease;
        text-align: center;
    }
    
    .mode-button:hover {
        background: rgba(79, 70, 229, 0.3);
        border-color: rgba(79, 70, 229, 0.7);
        transform: translateX(5px);
    }
    
    .mode-button.active {
        background: linear-gradient(135deg, rgba(79, 70, 229, 0.4) 0%, rgba(124, 58, 237, 0.4) 100%);
        border-color: rgba(79, 70, 229, 0.8);
    }
    
    h1, h2, h3 {
        color: #e0e7ff;
    }
    
    p, li, label {
        color: #cbd5e1;
    }
    
    .streamlit-expanderHeader {
        background: linear-gradient(135deg, rgba(79, 70, 229, 0.15) 0%, rgba(124, 58, 237, 0.15) 100%);
        color: #e0e7ff;
        border-radius: 10px;
        border: 1px solid rgba(79, 70, 229, 0.3);
        transition: all 0.3s ease;
    }
    
    .streamlit-expanderHeader:hover {
        background: linear-gradient(135deg, rgba(79, 70, 229, 0.25) 0%, rgba(124, 58, 237, 0.25) 100%);
        border-color: rgba(79, 70, 229, 0.5);
    }
    
    .success-message {
        background: linear-gradient(135deg, rgba(16, 185, 129, 0.15) 0%, rgba(5, 150, 105, 0.15) 100%);
        border-left: 4px solid #10b981;
        border-radius: 8px;
        padding: 15px;
        color: #6ee7b7;
        margin: 10px 0;
    }
    
    .info-box {
        background: rgba(59, 130, 246, 0.1);
        border-left: 4px solid #3b82f6;
        border-radius: 8px;
        padding: 15px;
        color: #93c5fd;
        margin: 10px 0;
    }
    
    code {
        background: rgba(79, 70, 229, 0.2);
        color: #c4b5fd;
        padding: 3px 8px;
        border-radius: 6px;
        border: 1px solid rgba(79, 70, 229, 0.3);
    }
    
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        background: rgba(26, 26, 63, 0.5);
        padding: 10px;
        border-radius: 12px;
    }
    
    .stTabs [data-baseweb="tab"] {
        background: rgba(79, 70, 229, 0.2);
        border-radius: 8px;
        color: #cbd5e1;
        border: 1px solid rgba(79, 70, 229, 0.3);
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
        color: white;
    }
</style>
""", unsafe_allow_html=True)

# Configuration
# Using ArborConfig from setup.py

# Enhanced Metadata Extractor for Claims Intelligence (Complete Version)
class EnhancedMetadataExtractor:
    """Advanced metadata extraction with multilingual support"""
    def __init__(self):
        # Claim number patterns - comprehensive
        self.claim_patterns = [
            r'claim\s*(?:number|no\.?|#|id|ref(?:erence)?)\s*:?\s*([A-Z0-9\-/]+)',
            r'(?:claim|policy)\s*([A-Z]{2,}\d{4,})',
            r'reference\s*(?:number|no\.?|#)?\s*:?\s*([A-Z0-9\-/]+)',
            r'\b([A-Z]{2,3}\-\d{6,})\b',  # Format: CLM-123456
            r'\b(CLM\d{6,})\b',
            r'\b(CLAIM\d{6,})\b'
        ]
        
        # Policy number patterns
        self.policy_patterns = [
            r'policy\s*(?:number|no\.?|#)?\s*:?\s*([A-Z0-9\-/]+)',
            r'\b(POL\d{6,})\b',
            r'\b([A-Z]{3}\d{7,})\b'
        ]
        
        # Date patterns - multiple formats
        self.date_patterns = [
            r'\b\d{1,2}[-/]\d{1,2}[-/]\d{2,4}\b',
            r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b',
            r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{4}\b',
            r'\b\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b'
        ]
        
        # Amount patterns - comprehensive
        self.amount_patterns = [
            r'\$\s?\d{1,3}(?:,\d{3})*(?:\.\d{2})?',
            r'USD\s?\d{1,3}(?:,\d{3})*(?:\.\d{2})?',
            r'EUR\s?\d{1,3}(?:,\d{3})*(?:\.\d{2})?',
            r'€\s?\d{1,3}(?:,\d{3})*(?:\.\d{2})?',
            r'\d{1,3}(?:,\d{3})*(?:\.\d{2})?\s?(?:dollars|USD|EUR|euros)',
            r'(?:amount|sum|total|payment|settlement)\s*:?\s*\$?\s?\d{1,3}(?:,\d{3})*(?:\.\d{2})?'
        ]
        
        # Urgency indicators - weighted by importance (multilingual)
        self.urgency_indicators = {
            # English
            'critical': 5, 'emergency': 5, 'urgent': 4, 'immediate': 4,
            'asap': 4, 'priority': 3, 'expedite': 3, 'time-sensitive': 3,
            'rush': 3, 'prompt': 2, 'quickly': 2, 'soon': 1,
            'deadline': 3, 'overdue': 4,
            # German
            'kritisch': 5, 'notfall': 5, 'dringend': 4, 'sofort': 4,
            'priorität': 3, 'eilt': 3, 'frist': 3, 'überfällig': 4,
            # French
            'critique': 5, 'urgence': 5, 'urgent': 4, 'immédiat': 4,
            # Spanish
            'crítico': 5, 'emergencia': 5, 'urgente': 4, 'inmediato': 4
        }
        
        # Document type keywords (multilingual)
        self.doc_type_keywords = {
            'claim': {
                'primary': ['claimant', 'incident', 'loss', 'damage', 'injury', 'accident',
                           'schadenfall', 'anspruch', 'réclamation', 'siniestro'],
                'secondary': ['claim', 'filed', 'reported', 'occurred', 'schaden']
            },
            'policy': {
                'primary': ['policy', 'coverage', 'insured', 'premium', 'deductible',
                           'versicherung', 'police', 'póliza'],
                'secondary': ['terms', 'conditions', 'benefits', 'exclusions', 'bedingungen']
            },
            'guideline': {
                'primary': ['procedure', 'guideline', 'process', 'standard', 'protocol',
                           'richtlinie', 'verfahren', 'procédure', 'procedimiento'],
                'secondary': ['step', 'instruction', 'requirement', 'must', 'should']
            },
            'regulation': {
                'primary': ['regulation', 'compliance', 'law', 'statute', 'requirement',
                           'vorschrift', 'gesetz', 'règlement', 'reglamento'],
                'secondary': ['legal', 'mandatory', 'regulatory', 'federal', 'state']
            }
        }
        
        # Status indicators
        self.status_patterns = {
            'approved': r'\b(?:approved|accepted|granted|genehmigt|approuvé|aprobado)\b',
            'denied': r'\b(?:denied|rejected|declined|disapproved|abgelehnt|refusé|denegado)\b',
            'pending': r'\b(?:pending|under review|in process|reviewing|ausstehend|en attente|pendiente)\b',
            'closed': r'\b(?:closed|settled|resolved|completed|geschlossen|fermé|cerrado)\b'
        }
        
        # Contact information patterns
        self.contact_patterns = {
            'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            'phone': r'\b(?:\+?1[-.]?)?\(?([0-9]{3})\)?[-.]?([0-9]{3})[-.]?([0-9]{4})\b'
        }
    
    def extract_claim_numbers(self, text: str) -> List[str]:
        """Extract all claim numbers with deduplication"""
        claim_numbers = set()
        for pattern in self.claim_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            claim_numbers.update(matches)
        return sorted(list(claim_numbers))
    
    def extract_policy_numbers(self, text: str) -> List[str]:
        """Extract policy numbers"""
        policy_numbers = set()
        for pattern in self.policy_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            policy_numbers.update(matches)
        return sorted(list(policy_numbers))
    
    def extract_dates(self, text: str) -> List[str]:
        """Extract all dates mentioned"""
        dates = []
        for pattern in self.date_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            dates.extend(matches)
        return dates[:10]
    
    def extract_amounts(self, text: str) -> List[Dict]:
        """Extract monetary amounts with context"""
        amounts = []
        for pattern in self.amount_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                amount_str = match.group(0)
                start = max(0, match.start() - 50)
                end = min(len(text), match.end() + 50)
                context = text[start:end].strip()
                amounts.append({'amount': amount_str, 'context': context})
        return amounts[:10]
    
    def calculate_urgency_score(self, text: str) -> Dict:
        """Calculate urgency score with detailed breakdown"""
        text_lower = text.lower()
        urgency_details = {
            'score': 0,
            'level': 'normal',
            'indicators_found': [],
            'total_mentions': 0
        }
        
        for indicator, weight in self.urgency_indicators.items():
            count = text_lower.count(indicator)
            if count > 0:
                urgency_details['score'] += count * weight
                urgency_details['total_mentions'] += count
                urgency_details['indicators_found'].append({
                    'term': indicator,
                    'count': count,
                    'weight': weight
                })
        
        if urgency_details['score'] >= 10:
            urgency_details['level'] = 'critical'
        elif urgency_details['score'] >= 6:
            urgency_details['level'] = 'high'
        elif urgency_details['score'] >= 3:
            urgency_details['level'] = 'medium'
        
        return urgency_details
    
    def detect_document_type(self, text: str, filename: str) -> Dict:
        """Enhanced document type detection with confidence scoring"""
        text_lower = text.lower()
        filename_lower = filename.lower()
        
        type_scores = {}
        for doc_type, keywords in self.doc_type_keywords.items():
            score = 0
            for keyword in keywords['primary']:
                score += text_lower.count(keyword) * 3
                if keyword in filename_lower:
                    score += 10
            for keyword in keywords['secondary']:
                score += text_lower.count(keyword) * 1
                if keyword in filename_lower:
                    score += 5
            type_scores[doc_type] = score
        
        if max(type_scores.values()) > 0:
            best_type = max(type_scores.items(), key=lambda x: x[1])
            confidence = min(best_type[1] / 20, 1.0)
            return {
                'type': best_type[0],
                'confidence': round(confidence, 2),
                'scores': type_scores
            }
        
        return {'type': 'general', 'confidence': 0.0, 'scores': type_scores}
    
    def extract_status(self, text: str) -> str:
        """Extract claim/document status"""
        text_lower = text.lower()
        for status, pattern in self.status_patterns.items():
            if re.search(pattern, text_lower):
                return status
        return 'unknown'
    
    def extract_contacts(self, text: str) -> Dict:
        """Extract contact information"""
        contacts = {'emails': [], 'phones': []}
        emails = re.findall(self.contact_patterns['email'], text)
        contacts['emails'] = list(set(emails))[:5]
        phones = re.findall(self.contact_patterns['phone'], text)
        contacts['phones'] = [f"({p[0]}) {p[1]}-{p[2]}" for p in phones[:5]]
        return contacts
    
    def extract_all_metadata(self, text: str, filename: str) -> Dict:
        """Comprehensive metadata extraction"""
        return {
            'claim_numbers': self.extract_claim_numbers(text),
            'policy_numbers': self.extract_policy_numbers(text),
            'dates': self.extract_dates(text),
            'amounts': self.extract_amounts(text),
            'urgency': self.calculate_urgency_score(text),
            'document_type': self.detect_document_type(text, filename),
            'status': self.extract_status(text),
            'contacts': self.extract_contacts(text),
            'word_count': len(text.split()),
            'char_count': len(text)
        }

# AI Summarization Engine
class SnippetGenerator:
    """Generate snippets and AI-powered summaries"""
    def __init__(self):
        self.summarizer = None
        if SEARCH_AVAILABLE:
            try:
                self.summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
            except Exception as e:
                self.summarizer = None
    
    def extract_snippets(self, content: str, query: str, max_snippets: int = 3) -> List[str]:
        """Extract relevant snippets from content"""
        if not SEARCH_AVAILABLE:
            words = query.lower().split()
            content_lower = content.lower()
            snippets = []
            for word in words:
                idx = content_lower.find(word)
                if idx != -1:
                    start = max(0, idx - 100)
                    end = min(len(content), idx + 100)
                    snippets.append(content[start:end].strip())
                    if len(snippets) >= max_snippets:
                        break
            return snippets
        
        # Advanced snippet extraction
        sentences = sent_tokenize(content)
        query_tokens = set(query.lower().split())
        
        sentence_scores = []
        for sentence in sentences:
            sentence_tokens = set(sentence.lower().split())
            score = len(query_tokens.intersection(sentence_tokens)) / len(query_tokens) if query_tokens else 0
            sentence_scores.append((score, sentence))
        
        sentence_scores.sort(key=lambda x: x[0], reverse=True)
        snippets = [sentence for score, sentence in sentence_scores[:max_snippets] if score > 0]
        return snippets
    
    def generate_summary(self, content: str, max_length: int = 150) -> str:
        """Generate AI-powered summary"""
        if not self.summarizer or len(content) < 100:
            if SEARCH_AVAILABLE:
                sentences = sent_tokenize(content)
                return ' '.join(sentences[:2]) if sentences else content[:max_length]
            else:
                return content[:max_length] + "..."
        
        try:
            if len(content) > 1024:
                content = content[:1024]
            summary = self.summarizer(content, max_length=max_length, min_length=30, do_sample=False)
            return summary[0]['summary_text']
        except Exception as e:
            if SEARCH_AVAILABLE:
                sentences = sent_tokenize(content)
                return ' '.join(sentences[:2]) if sentences else content[:max_length]
            else:
                return content[:max_length] + "..."

# Hybrid Search Engine with 65/35 BM25/Semantic ratio (Complete Implementation)
class HybridSearchEngine:
    """Advanced hybrid search with 65% BM25 + 35% Semantic + AI Summarization + Multilingual Support"""
    def __init__(self):
        self.documents = []
        self.metadata_extractor = EnhancedMetadataExtractor()
        self.snippet_generator = SnippetGenerator()
        self.bm25_weight = 0.65  # 65% weight for keyword search
        self.semantic_weight = 0.35  # 35% weight for semantic search
        
        if SEARCH_AVAILABLE:
            self.stemmer = PorterStemmer()
            # Multilingual stopwords support
            self.stop_words = set()
            for lang in ['english', 'german', 'french', 'spanish']:
                try:
                    self.stop_words.update(stopwords.words(lang))
                except:
                    pass
            
            # Initialize semantic model (multilingual)
            try:
                # Using LaBSE for better multilingual support
                self.semantic_model = SentenceTransformer('sentence-transformers/LaBSE')
                self.document_embeddings = None
                self.semantic_enabled = True
            except:
                try:
                    # Fallback to all-MiniLM
                    self.semantic_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')
                    self.document_embeddings = None
                    self.semantic_enabled = True
                except:
                    self.semantic_enabled = False
        else:
            self.semantic_enabled = False
        
        # BM25 parameters
        self.k1 = 1.2
        self.b = 0.75
        self.doc_frequencies = {}
        self.doc_lengths = []
        self.avg_doc_length = 0
        self.vocabulary = set()
        self.processed_docs = []
    
    def preprocess_text(self, text: str) -> List[str]:
        """Preprocess text with multilingual support"""
        if not SEARCH_AVAILABLE:
            return text.lower().split()
        
        text = re.sub(r'[^a-zA-ZäöüßÄÖÜàèéêëîïôûùÀÈÉÊËÎÏÔÛÙáéíóúñÁÉÍÓÚÑ\s]', ' ', text)
        tokens = word_tokenize(text.lower())
        tokens = [self.stemmer.stem(t) for t in tokens 
                 if t not in self.stop_words and len(t) > 2]
        return tokens
    
    def extract_text_from_file(self, file_path: Path) -> str:
        """Extract text from various file formats"""
        try:
            ext = file_path.suffix.lower()
            
            if ext == '.txt' or ext == '.md':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            elif ext == '.pdf' and SEARCH_AVAILABLE:
                try:
                    doc = fitz.open(str(file_path))
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    return text
                except:
                    return ""
            
            elif ext == '.docx' and SEARCH_AVAILABLE:
                try:
                    doc = DocxDocument(str(file_path))
                    return '\n'.join([p.text for p in doc.paragraphs])
                except:
                    return ""
            
            elif ext == '.pptx' and SEARCH_AVAILABLE:
                try:
                    prs = Presentation(str(file_path))
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except:
                    return ""
            
            else:
                return ""
        except Exception as e:
            return ""
    
    def index_documents(self, docs_path: Path):
        """Index all documents with both BM25 and semantic embeddings"""
        self.documents = []
        
        if not docs_path.exists():
            return
        
        supported_formats = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
        
        # Step 1: Load and process documents
        for file_path in docs_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_formats:
                text = self.extract_text_from_file(file_path)
                
                if text and len(text.strip()) > 50:
                    # Extract comprehensive metadata
                    metadata = self.metadata_extractor.extract_all_metadata(text, file_path.stem)
                    tokens = self.preprocess_text(text)
                    
                    self.documents.append({
                        'id': str(file_path),
                        'title': file_path.stem,
                        'content': text,
                        'file_path': str(file_path),
                        'file_type': file_path.suffix,
                        'tokens': tokens,
                        'metadata': metadata
                    })
        
        if not self.documents:
            return
        
        # Step 2: Build BM25 index
        self.processed_docs = [doc['tokens'] for doc in self.documents]
        self.vocabulary = set()
        for doc in self.processed_docs:
            self.vocabulary.update(doc)
        
        self.doc_frequencies = {}
        for token in self.vocabulary:
            self.doc_frequencies[token] = sum(1 for doc in self.processed_docs if token in doc)
        
        self.doc_lengths = [len(doc) for doc in self.processed_docs]
        self.avg_doc_length = sum(self.doc_lengths) / len(self.doc_lengths) if self.doc_lengths else 0
        
        # Step 3: Build semantic index with chunking
        if self.semantic_enabled:
            try:
                all_chunks = []
                chunk_to_doc = []
                
                for idx, doc in enumerate(self.documents):
                    if SEARCH_AVAILABLE:
                        sentences = sent_tokenize(doc['content'])
                        for sentence in sentences:
                            if len(sentence.strip()) > 20:
                                all_chunks.append(sentence.strip())
                                chunk_to_doc.append(idx)
                    else:
                        words = doc['content'].split()
                        for i in range(0, len(words), 100):
                            chunk = ' '.join(words[i:i+100])
                            if len(chunk) > 20:
                                all_chunks.append(chunk)
                                chunk_to_doc.append(idx)
                
                if all_chunks:
                    self.document_embeddings = self.semantic_model.encode(all_chunks, show_progress_bar=False)
                    self.chunk_to_doc_mapping = chunk_to_doc
                    
                    # Store document-level embeddings
                    for idx, doc in enumerate(self.documents):
                        doc_chunks = [i for i, d_idx in enumerate(chunk_to_doc) if d_idx == idx]
                        if doc_chunks:
                            doc['semantic_embedding'] = np.mean(self.document_embeddings[doc_chunks], axis=0)
            except Exception as e:
                self.semantic_enabled = False
    
    def calculate_bm25_score(self, query_tokens: List[str], doc_tokens: List[str], doc_length: int) -> float:
        """Calculate BM25 score for a document"""
        score = 0.0
        N = len(self.documents)
        
        for token in query_tokens:
            if token in doc_tokens:
                tf = doc_tokens.count(token)
                df = self.doc_frequencies.get(token, 0)
                if df > 0:
                    idf = math.log((N - df + 0.5) / (df + 0.5))
                    score += idf * (tf * (self.k1 + 1)) / (tf + self.k1 * (1 - self.b + self.b * doc_length / self.avg_doc_length))
        
        return score
    
    def calculate_semantic_score(self, query: str, doc: Dict) -> float:
        """Calculate semantic similarity score"""
        if not self.semantic_enabled or 'semantic_embedding' not in doc:
            return 0.0
        
        try:
            query_embedding = self.semantic_model.encode([query], show_progress_bar=False)
            similarity = cosine_similarity(query_embedding, [doc['semantic_embedding']])[0][0]
            return float(similarity)
        except:
            return 0.0
    
    def get_document_stats(self) -> Dict:
        """Get statistics about indexed documents"""
        stats = {
            'total_documents': len(self.documents),
            'by_type': {},
            'by_urgency': {'critical': 0, 'high': 0, 'medium': 0, 'normal': 0},
            'with_claim_numbers': 0,
            'with_amounts': 0,
            'with_contacts': 0
        }
        
        for doc in self.documents:
            metadata = doc.get('metadata', {})
            
            # Count by document type
            doc_type_info = metadata.get('document_type', {})
            doc_type = doc_type_info.get('type', 'general')
            stats['by_type'][doc_type] = stats['by_type'].get(doc_type, 0) + 1
            
            # Count by urgency level
            urgency_level = metadata.get('urgency', {}).get('level', 'normal')
            stats['by_urgency'][urgency_level] += 1
            
            # Count documents with specific metadata
            if metadata.get('claim_numbers'):
                stats['with_claim_numbers'] += 1
            if metadata.get('amounts'):
                stats['with_amounts'] += 1
            if metadata.get('contacts', {}).get('emails') or metadata.get('contacts', {}).get('phones'):
                stats['with_contacts'] += 1
        
        return stats
    
    def search(self, query: str, top_k: int = 10, 
               doc_type_filter: Optional[str] = None, 
               urgency_filter: Optional[str] = None,
               generate_summaries: bool = False) -> Dict:
        """Hybrid search with 65/35 BM25/Semantic ratio + AI summaries"""
        if not self.documents:
            return {
                'results': [],
                'total_results': 0,
                'bm25_weight': self.bm25_weight,
                'semantic_weight': self.semantic_weight,
                'semantic_enabled': self.semantic_enabled,
                'processing_time': 0
            }
        
        start_time = datetime.now()
        query_tokens = self.preprocess_text(query)
        scored_docs = []
        
        # Calculate scores for each document
        for i, doc in enumerate(self.documents):
            # Apply filters
            metadata = doc.get('metadata', {})
            doc_type_info = metadata.get('document_type', {})
            
            if doc_type_filter and doc_type_info.get('type') != doc_type_filter:
                continue
            
            if urgency_filter and metadata.get('urgency', {}).get('level') != urgency_filter:
                continue
            
            # BM25 score
            bm25_score = self.calculate_bm25_score(query_tokens, self.processed_docs[i], self.doc_lengths[i])
            
            # Semantic score
            semantic_score = self.calculate_semantic_score(query, doc) if self.semantic_enabled else 0.0
            
            scored_docs.append({
                'doc': doc,
                'bm25_score': bm25_score,
                'semantic_score': semantic_score
            })
        
        # Normalize scores
        if scored_docs:
            max_bm25 = max([d['bm25_score'] for d in scored_docs], default=1)
            max_semantic = max([d['semantic_score'] for d in scored_docs], default=1)
            
            for item in scored_docs:
                norm_bm25 = item['bm25_score'] / max_bm25 if max_bm25 > 0 else 0
                norm_semantic = item['semantic_score'] / max_semantic if max_semantic > 0 else 0
                
                # Hybrid score with 65/35 ratio
                item['combined_score'] = (self.bm25_weight * norm_bm25) + (self.semantic_weight * norm_semantic)
                item['norm_bm25'] = norm_bm25
                item['norm_semantic'] = norm_semantic
        
        # Sort by combined score
        scored_docs.sort(key=lambda x: x['combined_score'], reverse=True)
        
        # Prepare results with snippets and summaries
        results = []
        for item in scored_docs[:top_k]:
            doc = item['doc']
            result = {
                **doc,
                'bm25_score': item['bm25_score'],
                'semantic_score': item['semantic_score'],
                'combined_score': item['combined_score'],
                'norm_bm25': round(item['norm_bm25'], 3),
                'norm_semantic': round(item['norm_semantic'], 3),
                'snippets': self.snippet_generator.extract_snippets(doc['content'], query)
            }
            
            # Generate AI summary if requested
            if generate_summaries:
                result['summary'] = self.snippet_generator.generate_summary(doc['content'])
            
            # Extract relevant chunks for display
            if self.semantic_enabled and 'semantic_embedding' in doc:
                # Find most relevant chunks
                try:
                    doc_idx = self.documents.index(doc)
                    doc_chunks_indices = [i for i, d_idx in enumerate(self.chunk_to_doc_mapping) if d_idx == doc_idx]
                    
                    if doc_chunks_indices:
                        query_embedding = self.semantic_model.encode([query], show_progress_bar=False)
                        chunk_embeddings = self.document_embeddings[doc_chunks_indices]
                        similarities = cosine_similarity(query_embedding, chunk_embeddings)[0]
                        
                        top_chunk_indices = np.argsort(similarities)[-3:][::-1]
                        relevant_chunks = []
                        for idx in top_chunk_indices:
                            actual_idx = doc_chunks_indices[idx]
                            if actual_idx < len(self.chunk_to_doc_mapping):
                                # Get chunk text from document
                                if SEARCH_AVAILABLE:
                                    sentences = sent_tokenize(doc['content'])
                                    if idx < len(sentences):
                                        relevant_chunks.append(sentences[idx])
                        
                        if relevant_chunks:
                            result['relevant_chunks'] = relevant_chunks
                except:
                    pass
            
            results.append(result)
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return {
            'query': query,
            'results': results,
            'total_results': len(results),
            'bm25_weight': self.bm25_weight,
            'semantic_weight': self.semantic_weight,
            'semantic_enabled': self.semantic_enabled,
            'processing_time': round(processing_time, 3),
            'filter_applied': {
                'doc_type': doc_type_filter,
                'urgency': urgency_filter
            }
        }

# Initialize session state
if 'initialized' not in st.session_state:
    st.session_state.initialized = ARBOR_AVAILABLE
    st.session_state.messages = []
    st.session_state.mode = "ask"
    # Initialize with HybridSearchEngine
    if SEARCH_AVAILABLE:
        st.session_state.searcher = HybridSearchEngine()
    else:
        st.session_state.searcher = None
    st.session_state.search_indexed = False
    st.session_state.stats = {
        'total_documents': 0,
        'vectorstore_status': 'Not initialized',
        'calendar_status': 'Not connected',
        'sheets_status': 'Not connected',
        'email_status': 'Not connected',
        'search_documents': 0,
        'last_refresh': None
    }
    
    # Update stats if Arbor is available
    if ARBOR_AVAILABLE:
        try:
            st.session_state.stats['total_documents'] = len(all_documents) if all_documents else 0
            st.session_state.stats['vectorstore_status'] = 'Active' if vectorstore else 'Inactive'
            st.session_state.stats['calendar_status'] = 'Connected' if calendar_connector.service else 'Not connected'
            st.session_state.stats['sheets_status'] = 'Connected' if sheets_connector.service else 'Not connected'
            st.session_state.stats['email_status'] = 'Ready'
        except:
            pass

# Sidebar
with st.sidebar:
    st.markdown("### Control Panel")
    st.markdown("---")
    
    # Mode selector
    st.markdown("#### Select Mode")
    
    if st.button("Ask Arbor", use_container_width=True, 
                 type="primary" if st.session_state.mode == "ask" else "secondary"):
        st.session_state.mode = "ask"
        st.rerun()
    
    if st.button("Tell Arbor", use_container_width=True,
                 type="primary" if st.session_state.mode == "do" else "secondary"):
        st.session_state.mode = "do"
        st.rerun()
    
    st.markdown("---")
    
    # System controls
    st.markdown("#### System Controls")
    
    corpus_path = st.text_input(
        "Document Repository:",
        value="./arbor_data/documents",
        help="Path to your documents folder"
    )
    
    if st.button("Initialize System", use_container_width=True):
        if not ARBOR_AVAILABLE:
            st.error("setup.py not available. Please ensure setup.py is in the same directory.")
        else:
            with st.spinner("Initializing Arbor Enhanced..."):
                try:
                    # System is already initialized from setup.py
                    # Just update stats
                    st.session_state.stats['total_documents'] = len(all_documents) if all_documents else 0
                    st.session_state.stats['vectorstore_status'] = 'Active' if vectorstore else 'Inactive'
                    st.session_state.stats['calendar_status'] = 'Connected' if calendar_connector.service else 'Not connected'
                    st.session_state.stats['sheets_status'] = 'Connected' if sheets_connector.service else 'Not connected'
                    st.session_state.stats['email_status'] = 'Ready'
                    st.session_state.initialized = True
                    st.success("System initialized!")
                    st.balloons()
                    st.rerun()
                except Exception as e:
                    st.error(f"Initialization failed: {e}")
    
    if st.button("Refresh Data", use_container_width=True):
        if not ARBOR_AVAILABLE:
            st.error("setup.py not available.")
        else:
            with st.spinner("Refreshing data sources..."):
                try:
                    result = refresh_all_data()
                    st.session_state.stats['last_refresh'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    st.session_state.stats['total_documents'] = len(all_documents) if all_documents else 0
                    st.success(f"{result}")
                except Exception as e:
                    st.error(f"Refresh failed: {e}")
    
    if st.button("Clear Chat", use_container_width=True):
        st.session_state.messages = []
        st.rerun()
    
    st.markdown("---")
    
    # What Does This Do section
    with st.expander("What Does This Do?", expanded=False):
        st.markdown("""
        ### Arbor : Your AI Assistant
        
        **Arbor** is your intelligent enterprise AI assistant that seamlessly integrates with:
        
        #### Integrations
        - **Google Drive**: Document search & analysis
        - **Google Calendar**: Meeting intelligence & scheduling
        - **Google Sheets**: Task & data management
        - **Gmail**: Automated email reports
        
        #### Two Modes
        
        **Ask Arbor (Q&A Mode)**
        - Ask questions about your documents
        - Get intelligent insights
        - Bilingual support (English/German)
        - Powered by RAG technology
        
        **Tell Arbor (Action Mode)**
        - Execute autonomous tasks
        - Generate reports & analytics
        - Manage calendar & tasks
        - Send automated emails
        
        #### Key Features
        - Advanced analytics & reporting
        - Calendar intelligence with free slot detection
        - Document intelligence & action extraction
        - AI-powered task recommendations
        - Productivity pattern analysis
        - Email automation with filters
        
        #### How to Use
        
        1. **Initialize**: Click "Initialize System"
        2. **Choose Mode**: Select Ask or Tell
        3. **Interact**: Type your query/command
        4. **Review**: Check dashboard for stats
        
        #### Example Commands
        
        **Ask Mode:**
        - "What are my high priority tasks?"
        - "Was sind meine Aufgaben?" (German)
        - "What should I focus on today?"
        
        **Tell Mode:**
        - "Generate weekly task report and email to user@example.com"
        - "Find free time slots for 2-hour meeting this week"
        - "Email overdue tasks to manager@example.com"
        - "Suggest optimal task order for today"
        
        #### Privacy
        All processing is done locally using your credentials. Data is never shared externally.
        """)
    
    st.markdown("---")
    
    # System info
    st.markdown("#### System Info")
    st.markdown(f"""
    <div style='font-size: 0.85rem; color: #94a3b8;'>
        <strong>Model:</strong> DeepSeek R1 7B<br>
        <strong>Status:</strong> {'Ready' if st.session_state.initialized else 'Not Initialized'}<br>
        <strong>Backend:</strong> {'Connected' if ARBOR_AVAILABLE else 'Not Available'}
    </div>
    """, unsafe_allow_html=True)

# Main content
col1, col2, col3 = st.columns([1, 3, 1])
with col2:
    st.markdown("""
    <div class="main-header">
        <h1>Arbor</h1>
        <p>Your Intelligent Enterprise AI Assistant</p>
    </div>
    """, unsafe_allow_html=True)

# Tabs for different sections
tab1, tab2, tab3, tab4 = st.tabs(["Chat", "Search Documents", "Dashboard", "Examples"])

with tab1:
    # Mode indicator
    mode_emoji = " " if st.session_state.mode == "ask" else " "
    mode_text = "Ask Arbor" if st.session_state.mode == "ask" else "Tell Arbor"
    mode_color = "#3b82f6" if st.session_state.mode == "ask" else "#8b5cf6"
    
    st.markdown(f"""
    <div style='background: linear-gradient(135deg, {mode_color}20 0%, {mode_color}10 100%); 
                padding: 15px; border-radius: 12px; border-left: 4px solid {mode_color}; margin-bottom: 20px;'>
        <h3 style='margin: 0; color: #e0e7ff;'>{mode_emoji} {mode_text} Mode</h3>
        <p style='margin: 5px 0 0 0; color: #cbd5e1; font-size: 0.9rem;'>
            {'Ask questions and get intelligent answers' if st.session_state.mode == 'ask' else 'Execute autonomous tasks and actions'}
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Chat display
    chat_container = st.container()
    
    with chat_container:
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
                
                if message["role"] == "assistant" and "details" in message:
                    with st.expander("Details"):
                        st.json(message["details"])
    
    # Chat input
    if prompt := st.chat_input(
        "Ask a question..." if st.session_state.mode == "ask" else "Tell Arbor what to do...",
        key="chat_input"
    ):
        if not st.session_state.initialized or not ARBOR_AVAILABLE:
            st.warning("Please initialize the system first using the sidebar.")
        else:
            # Add user message
            st.session_state.messages.append({"role": "user", "content": prompt})
            
            # Display user message
            with st.chat_message("user"):
                st.markdown(prompt)
            
            # Generate response
            with st.chat_message("assistant"):
                with st.spinner("Processing..."):
                    try:
                        if st.session_state.mode == "ask":
                            # Ask Arbor mode
                            result = ask_arbor(prompt, verbose=False)
                            
                            if 'error' in result:
                                response = f"**Error:** {result['error']}"
                                details = {"mode": "ask", "query": prompt, "error": result['error']}
                            else:
                                answer = result.get('answer', 'No answer available')
                                response = f"**Arbor's Response:**\n\n{answer}"
                                details = {
                                    "mode": "ask",
                                    "query": prompt,
                                    "timestamp": result.get('timestamp', datetime.now().isoformat())
                                }
                        else:
                            # Tell Arbor mode
                            result = arbor_do(prompt)
                            
                            response = f"**Action Completed:**\n\n{result}"
                            details = {
                                "mode": "do",
                                "command": prompt,
                                "timestamp": datetime.now().isoformat()
                            }
                        
                        st.markdown(response)
                        
                        # Show details in expander
                        with st.expander("Details"):
                            st.json(details)
                        
                        # Add assistant message
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": response,
                            "details": details
                        })
                        
                    except Exception as e:
                        error_msg = f"**Error:** {str(e)}"
                        st.error(error_msg)
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": error_msg,
                            "details": {"error": str(e)}
                        })

with tab2:
    st.markdown("### Hybrid Search Engine: 65% BM25 + 35% Semantic AI")
    st.markdown("**Multilingual Support:** English, German, French, Spanish | **AI-Powered Summaries** | **Advanced Metadata Extraction**")
    
    if not SEARCH_AVAILABLE:
        st.error(f"""
        **Search functionality requires additional packages**
        
        Please install: `pip install nltk PyMuPDF python-docx python-pptx sentence-transformers scikit-learn transformers torch`
        
        Error: {IMPORT_ERROR if 'IMPORT_ERROR' in globals() else 'Import failed'}
        """)
    else:
        # Index documents button with stats
        if st.session_state.search_indexed and st.session_state.searcher:
            stats = st.session_state.searcher.get_document_stats()
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Documents", stats['total_documents'])
            with col2:
                st.metric("With Claims", stats['with_claim_numbers'])
            with col3:
                st.metric("With Amounts", stats['with_amounts'])
            with col4:
                semantic_status = "ON" if st.session_state.searcher.semantic_enabled else "OFF"
                st.metric("Semantic AI", semantic_status)
            
            # Document type breakdown
            if stats['by_type']:
                st.markdown("**By Document Type:**")
                type_cols = st.columns(len(stats['by_type']))
                for i, (doc_type, count) in enumerate(stats['by_type'].items()):
                    with type_cols[i]:
                        st.metric(doc_type.title(), count)
            
            # Urgency breakdown
            st.markdown("**By Urgency Level:**")
            urg_col1, urg_col2, urg_col3, urg_col4 = st.columns(4)
            with urg_col1:
                st.metric("Critical", stats['by_urgency']['critical'], 
                         delta="Highest Priority" if stats['by_urgency']['critical'] > 0 else None,
                         delta_color="inverse")
            with urg_col2:
                st.metric("High", stats['by_urgency']['high'])
            with urg_col3:
                st.metric("Medium", stats['by_urgency']['medium'])
            with urg_col4:
                st.metric("Normal", stats['by_urgency']['normal'])
            
            if st.button("Re-index Documents", use_container_width=True):
                with st.spinner("Re-indexing documents..."):
                    st.session_state.searcher = HybridSearchEngine()
                    docs_path = Path(ArborConfig.DOCS_DIR)
                    st.session_state.searcher.index_documents(docs_path)
                    st.session_state.stats['search_documents'] = len(st.session_state.searcher.documents)
                    st.success("Documents re-indexed!")
                    st.rerun()
        else:
            col1, col2 = st.columns([3, 1])
            with col1:
                st.info("**Hybrid search engine** with 65% keyword matching + 35% AI semantic understanding. Supports 100+ languages including English, German, French, and Spanish.")
            with col2:
                if st.button("Index Docs", use_container_width=True, type="primary"):
                    if ARBOR_AVAILABLE:
                        with st.spinner("Indexing documents with hybrid search..."):
                            st.session_state.searcher = HybridSearchEngine()
                            docs_path = Path(ArborConfig.DOCS_DIR)
                            st.session_state.searcher.index_documents(docs_path)
                            st.session_state.search_indexed = True
                            st.session_state.stats['search_documents'] = len(st.session_state.searcher.documents)
                            
                            semantic_status = "Enabled" if st.session_state.searcher.semantic_enabled else "Disabled"
                            st.success(f"Indexed {st.session_state.stats['search_documents']} documents! Semantic: {semantic_status}")
                            st.rerun()
                    else:
                        st.error("Arbor not initialized")
        
        if not st.session_state.search_indexed:
            st.markdown("""
            <div style='text-align: center; padding: 3rem 0;'>
                <h3 style='color: #a78bfa;'>Index Your Documents</h3>
                <p style='color: #cbd5e1; font-size: 1.1rem;'>
                    Click "Index Docs" to enable advanced hybrid search
                </p>
                <div style='background: rgba(79, 70, 229, 0.1); padding: 2rem; border-radius: 15px; margin: 2rem 0;'>
                    <h4 style='color: #c4b5fd;'>Features</h4>
                    <ul style='color: #cbd5e1; text-align: left; display: inline-block;'>
                        <li>65/35 BM25/Semantic Hybrid Search</li>
                        <li>Multilingual: EN, DE, FR, ES support</li>
                        <li>AI-Powered Summarization</li>
                        <li>Advanced Metadata Extraction</li>
                        <li>Urgency Detection & Classification</li>
                        <li>PDF, DOCX, PPTX, TXT, MD support</li>
                    </ul>
                </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown("---")
            
            # Search interface
            search_col1, search_col2 = st.columns([4, 1])
            with search_col1:
                search_query = st.text_input(
                    "Search documents (100+ languages supported)",
                    placeholder="e.g., urgent medical claims, dringende Anträge, réclamations urgentes...",
                    key="hybrid_doc_search_input",
                    help="Supports English, German, French, Spanish and 100+ other languages"
                )
            with search_col2:
                st.markdown("<br>", unsafe_allow_html=True)
                search_btn = st.button("Search", type="primary", use_container_width=True, key="hybrid_search_btn")
            
            # Advanced Filters
            st.markdown("**Advanced Filters:**")
            filter_col1, filter_col2, filter_col3, filter_col4 = st.columns(4)
            
            with filter_col1:
                doc_type_filter = st.selectbox(
                    "Document Type",
                    options=["All", "claim", "policy", "guideline", "regulation", "general"],
                    key="hybrid_doc_type_filter",
                    help="Filter by document classification"
                )
            
            with filter_col2:
                urgency_filter = st.selectbox(
                    "Urgency Level",
                    options=["All", "critical", "high", "medium", "normal"],
                    key="hybrid_urgency_filter",
                    help="Filter by urgency/priority level"
                )
            
            with filter_col3:
                top_k = st.slider("Results", 5, 20, 10, key="hybrid_search_top_k",
                                 help="Number of results to display")
            
            with filter_col4:
                generate_summaries = st.checkbox("AI Summaries", value=False, key="hybrid_gen_summaries",
                                                help="Generate AI-powered document summaries (slower)")
            
            # Perform search
            if search_query and (search_btn or search_query):
                doc_filter = None if doc_type_filter == "All" else doc_type_filter
                urg_filter = None if urgency_filter == "All" else urgency_filter
                
                with st.spinner("Searching with hybrid algorithm..."):
                    search_results = st.session_state.searcher.search(
                        query=search_query,
                        top_k=top_k,
                        doc_type_filter=doc_filter,
                        urgency_filter=urg_filter,
                        generate_summaries=generate_summaries
                    )
                
                results = search_results['results']
                processing_time = search_results['processing_time']
                
                # Search Insights
                st.markdown("<br>", unsafe_allow_html=True)
                insight_col1, insight_col2, insight_col3, insight_col4 = st.columns(4)
                
                with insight_col1:
                    st.markdown(f"""
                    <div class="stat-card">
                        <div class="stat-number">{len(results)}</div>
                        <div class="stat-label">Results Found</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with insight_col2:
                    st.markdown(f"""
                    <div class="stat-card">
                        <div class="stat-number">{processing_time}s</div>
                        <div class="stat-label">Processing Time</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with insight_col3:
                    st.markdown(f"""
                    <div class="stat-card">
                        <div class="stat-number">65/35</div>
                        <div class="stat-label">BM25/Semantic</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with insight_col4:
                    filter_text = f"{doc_filter or 'All'}"
                    if urg_filter:
                        filter_text += f" / {urg_filter}"
                    st.markdown(f"""
                    <div class="stat-card">
                        <div class="stat-number" style="font-size: 1.2rem;">{filter_text}</div>
                        <div class="stat-label">Filters Applied</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                if not results:
                    st.info("No documents found matching your query. Try different keywords, adjust filters, or search in another language.")
                else:
                    # Results Header
                    st.markdown(f"""
                    <div style='background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
                                padding: 1.5rem; border-radius: 15px; color: white; margin: 2rem 0;'>
                        <h3 style='margin: 0;'>Hybrid Search Results</h3>
                        <p style='margin: 0.5rem 0 0 0; opacity: 0.9;'>
                            Found {len(results)} documents for "{search_query}" using 65% keyword + 35% semantic matching
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Display results
                    for i, result in enumerate(results):
                        metadata = result['metadata']
                        urgency = metadata['urgency']
                        doc_type_info = metadata.get('document_type', {})
                        
                        # Urgency badge color
                        urgency_colors = {
                            'critical': ('#fee2e2', '#991b1b'),
                            'high': ('#fed7aa', '#9a3412'),
                            'medium': ('#fef08a', '#854d0e'),
                            'normal': ('#e0e7ff', '#3730a3')
                        }
                        urg_bg, urg_text = urgency_colors.get(urgency['level'], ('#e0e7ff', '#3730a3'))
                        
                        # Document type color
                        type_colors = {
                            'claim': ('#dbeafe', '#1e40af'),
                            'policy': ('#d1fae5', '#065f46'),
                            'guideline': ('#fef3c7', '#92400e'),
                            'regulation': ('#fce7f3', '#831843'),
                            'general': ('#e0e7ff', '#3730a3')
                        }
                        type_bg, type_text = type_colors.get(doc_type_info.get('type', 'general'), ('#e0e7ff', '#3730a3'))
                        
                        st.markdown(f"""
                        <div style='background: rgba(79, 70, 229, 0.05); border: 1px solid rgba(79, 70, 229, 0.2);
                                    border-radius: 15px; padding: 1.5rem; margin: 1rem 0;
                                    transition: all 0.3s ease;'>
                            <div style='display: flex; justify-content: space-between; align-items: start;'>
                                <div style='flex: 1;'>
                                    <div style='margin-bottom: 0.5rem;'>
                                        <span style='background: {urg_bg}; color: {urg_text}; 
                                                     padding: 0.4rem 1rem; border-radius: 20px;
                                                     font-size: 0.85rem; font-weight: 600;'>
                                            {urgency['level'].upper()}
                                        </span>
                                        <span style='background: {type_bg}; color: {type_text}; 
                                                     padding: 0.4rem 1rem; border-radius: 20px;
                                                     font-size: 0.85rem; font-weight: 600; margin-left: 0.5rem;'>
                                            {doc_type_info.get('type', 'general').upper()}
                                        </span>
                                    </div>
                                    <h3 style='color: #e0e7ff; margin: 0.5rem 0;'>
                                        {i+1}. {result['title']}
                                    </h3>
                                    <p style='color: #94a3b8; font-size: 0.9rem; margin: 0.3rem 0;'>
                                        {result['file_type']} • {metadata['word_count']} words • 
                                        Type Confidence: {doc_type_info.get('confidence', 0):.0%}
                                    </p>
                                </div>
                            </div>
                        """, unsafe_allow_html=True)
                        
                        # Score breakdown
                        score_col1, score_col2, score_col3 = st.columns(3)
                        
                        with score_col1:
                            st.markdown(f"""
                            <div style='text-align: center; padding: 1rem; background: rgba(79, 70, 229, 0.15);
                                        border-radius: 10px; border: 2px solid #4f46e5;'>
                                <div style='font-size: 1.5rem; font-weight: 700; color: #6366f1;'>
                                    {result['norm_bm25']:.3f}
                                </div>
                                <div style='font-size: 0.75rem; color: #cbd5e1;'>Keyword (65%)</div>
                            </div>
                            """, unsafe_allow_html=True)
                        
                        with score_col2:
                            st.markdown(f"""
                            <div style='text-align: center; padding: 1rem; background: rgba(124, 58, 237, 0.15);
                                        border-radius: 10px; border: 2px solid #7c3aed;'>
                                <div style='font-size: 1.5rem; font-weight: 700; color: #a78bfa;'>
                                    {result['norm_semantic']:.3f}
                                </div>
                                <div style='font-size: 0.75rem; color: #cbd5e1;'>Semantic (35%)</div>
                            </div>
                            """, unsafe_allow_html=True)
                        
                        with score_col3:
                            st.markdown(f"""
                            <div style='text-align: center; padding: 1rem; background: rgba(16, 185, 129, 0.15);
                                        border-radius: 10px; border: 2px solid #10b981;'>
                                <div style='font-size: 1.5rem; font-weight: 700; color: #10b981;'>
                                    {result['combined_score']:.3f}
                                </div>
                                <div style='font-size: 0.75rem; color: #cbd5e1;'>Final Score</div>
                            </div>
                            """, unsafe_allow_html=True)
                        
                        # Enhanced Metadata Display
                        st.markdown(f"""
                        <div style='background: rgba(79, 70, 229, 0.08); padding: 1rem; border-radius: 10px; margin: 1rem 0;
                                    border-left: 4px solid #7c3aed;'>
                            <strong style='color: #c4b5fd;'>Extracted Metadata</strong>
                        """, unsafe_allow_html=True)
                        
                        meta_col1, meta_col2 = st.columns(2)
                        
                        with meta_col1:
                            if metadata.get('claim_numbers'):
                                st.markdown(f"**Claim Numbers:** {', '.join(metadata['claim_numbers'][:3])}")
                            
                            if metadata.get('policy_numbers'):
                                st.markdown(f"**Policy Numbers:** {', '.join(metadata['policy_numbers'][:3])}")
                            
                            if metadata.get('dates'):
                                st.markdown(f"**Dates:** {', '.join(metadata['dates'][:3])}")
                            
                            if metadata.get('status') and metadata['status'] != 'unknown':
                                st.markdown(f"**Status:** {metadata['status'].title()}")
                        
                        with meta_col2:
                            if metadata.get('amounts'):
                                amounts_str = ", ".join([amt['amount'] for amt in metadata['amounts'][:3]])
                                st.markdown(f"**Amounts:** {amounts_str}")
                            
                            contacts = metadata.get('contacts', {})
                            if contacts.get('emails'):
                                st.markdown(f"**Emails:** {', '.join(contacts['emails'][:2])}")
                            
                            if contacts.get('phones'):
                                st.markdown(f"**Phones:** {', '.join(contacts['phones'][:2])}")
                        
                        # Urgency breakdown
                        if urgency.get('indicators_found'):
                            with st.expander("Urgency Analysis Details"):
                                st.markdown(f"**Urgency Score:** {urgency.get('score', 0)}")
                                st.markdown(f"**Total Mentions:** {urgency.get('total_mentions', 0)}")
                                st.markdown("**Indicators Found:**")
                                for indicator in urgency['indicators_found']:
                                    st.markdown(f"- **{indicator['term']}**: {indicator['count']} times (weight: {indicator['weight']})")
                        
                        st.markdown("</div>", unsafe_allow_html=True)
                        
                        # Key Excerpts
                        if result.get('snippets'):
                            st.markdown("**Key Excerpts:**")
                            for snippet in result['snippets']:
                                st.markdown(f"""
                                <div style='background: rgba(245, 158, 11, 0.1); padding: 1rem; 
                                            border-radius: 10px; margin: 0.5rem 0; 
                                            border-left: 4px solid #f59e0b; color: #cbd5e1;
                                            font-style: italic;'>
                                    {snippet}
                                </div>
                                """, unsafe_allow_html=True)
                        
                        # Relevant Chunks
                        if result.get('relevant_chunks'):
                            with st.expander("Semantically Relevant Sections"):
                                for chunk in result['relevant_chunks']:
                                    st.markdown(f"""
                                    <div style='background: rgba(124, 58, 237, 0.1); padding: 1rem; 
                                                border-radius: 10px; margin: 0.5rem 0; 
                                                border-left: 3px solid #7c3aed; color: #cbd5e1;'>
                                        {chunk}
                                    </div>
                                    """, unsafe_allow_html=True)
                        
                        # AI Summary
                        if result.get('summary'):
                            with st.expander("AI-Generated Summary"):
                                st.markdown(f"""
                                <div style='background: linear-gradient(135deg, #dbeafe 0%, #bfdbfe 100%);
                                            padding: 1.5rem; border-radius: 10px; color: #1e40af;
                                            line-height: 1.8;'>
                                    {result['summary']}
                                </div>
                                """, unsafe_allow_html=True)
                        
                        # Document preview
                        with st.expander("Document Preview"):
                            preview_text = result['content'][:1500] + "..." if len(result['content']) > 1500 else result['content']
                            st.text_area("", preview_text, height=200, disabled=True, key=f"hybrid_search_preview_{i}")
                        
                        # File details
                        with st.expander("File Details"):
                            st.text(f"Path: {result['file_path']}")
                            if st.button(f"Show Folder Path", key=f"hybrid_show_folder_{i}"):
                                folder_path = Path(result['file_path']).parent
                                st.code(str(folder_path), language=None)
                        
                        st.markdown("</div>", unsafe_allow_html=True)

with tab3:
    st.markdown("### System Dashboard")
    
    # Stats grid
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.markdown(f"""
        <div class="stat-card">
            <div class="stat-number">{st.session_state.stats['total_documents']}</div>
            <div class="stat-label">Documents Indexed (RAG)</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div class="stat-card">
            <div class="stat-number">{st.session_state.stats.get('search_documents', 0)}</div>
            <div class="stat-label">Searchable Documents</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        status_emoji = "OK" if st.session_state.initialized else "NOT OK"
        st.markdown(f"""
        <div class="stat-card">
            <div class="stat-number">{status_emoji}</div>
            <div class="stat-label">System Status</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        chat_count = len(st.session_state.messages) // 2
        st.markdown(f"""
        <div class="stat-card">
            <div class="stat-number">{chat_count}</div>
            <div class="stat-label">Conversations</div>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Component status
    st.markdown("#### Component Status")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown(f"""
        <div class="info-box">
            <strong>Vector Store:</strong> {st.session_state.stats['vectorstore_status']}<br>
            <strong>Calendar:</strong> {st.session_state.stats['calendar_status']}<br>
            <strong>Search Engine:</strong> {'Ready' if st.session_state.search_indexed else 'Not indexed'}
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div class="info-box">
            <strong>Sheets:</strong> {st.session_state.stats['sheets_status']}<br>
            <strong>Email:</strong> {st.session_state.stats.get('email_status', 'Not connected')}<br>
            <strong>Last Refresh:</strong> {st.session_state.stats['last_refresh'] or 'Never'}
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Recent activity
    st.markdown("#### Recent Activity")
    
    if st.session_state.messages:
        recent_df = pd.DataFrame([
            {
                "Time": datetime.now().strftime("%H:%M:%S"),
                "Mode": msg.get("details", {}).get("mode", "unknown").upper() if msg["role"] == "assistant" else "USER",
                "Message": msg["content"][:50] + "..." if len(msg["content"]) > 50 else msg["content"]
            }
            for msg in st.session_state.messages[-5:]
        ])
        st.dataframe(recent_df, use_container_width=True)
    else:
        st.info("No recent activity. Start chatting to see activity here!")

with tab4:
    st.markdown("### Example Queries")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("""
        <div style='background: rgba(79, 70, 229, 0.1); padding: 20px; border-radius: 12px; 
                    border: 1px solid rgba(79, 70, 229, 0.3);'>
            <h4 style='color: #a78bfa; margin-bottom: 15px;'>Ask Mode Examples</h4>
            <ul style='line-height: 2; color: #cbd5e1;'>
                <li>What are my high priority tasks?</li>
                <li>Show me meetings for today</li>
                <li>Was sind meine Aufgaben? (German)</li>
                <li>What should I focus on today?</li>
                <li>Find documents about project X</li>
                <li>When is my next deadline?</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown("""
        <div style='background: rgba(124, 58, 237, 0.1); padding: 20px; border-radius: 12px; 
                    border: 1px solid rgba(124, 58, 237, 0.3);'>
            <h4 style='color: #c4b5fd; margin-bottom: 15px;'>Tell Mode Examples</h4>
            <ul style='line-height: 2; color: #cbd5e1;'>
                <li>Generate weekly task report and email to user@example.com</li>
                <li>Find free time slots for 2-hour meeting this week</li>
                <li>Email overdue tasks to manager@example.com</li>
                <li>Suggest optimal task order for today</li>
                <li>Analyze my productivity patterns</li>
                <li>Extract action items from meeting notes</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Quick actions
    st.markdown("### Quick Actions")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("Today's Schedule", use_container_width=True):
            if ARBOR_AVAILABLE:
                st.session_state.mode = "ask"
                prompt = "Show me today's schedule"
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                try:
                    result = ask_arbor(prompt, verbose=False)
                    if 'error' in result:
                        response = f"**Error:** {result['error']}"
                        details = {"mode": "ask", "query": prompt, "error": result['error']}
                    else:
                        answer = result.get('answer', 'No answer available')
                        response = f"**Arbor's Response:**\n\n{answer}"
                        details = {"mode": "ask", "query": prompt, "timestamp": result.get('timestamp', datetime.now().isoformat())}
                    
                    st.session_state.messages.append({"role": "assistant", "content": response, "details": details})
                except Exception as e:
                    error_msg = f"**Error:** {str(e)}"
                    st.session_state.messages.append({"role": "assistant", "content": error_msg, "details": {"error": str(e)}})
                
                st.rerun()
            else:
                st.error("Arbor not available")
    
    with col2:
        if st.button("Weekly Report", use_container_width=True):
            if ARBOR_AVAILABLE:
                st.session_state.mode = "do"
                prompt = "Generate weekly task report"
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                try:
                    result = arbor_do(prompt)
                    response = f"**Action Completed:**\n\n{result}"
                    details = {"mode": "do", "command": prompt, "timestamp": datetime.now().isoformat()}
                    st.session_state.messages.append({"role": "assistant", "content": response, "details": details})
                except Exception as e:
                    error_msg = f"**Error:** {str(e)}"
                    st.session_state.messages.append({"role": "assistant", "content": error_msg, "details": {"error": str(e)}})
                
                st.rerun()
            else:
                st.error("Arbor not available")
    
    with col3:
        if st.button("Find Free Slots", use_container_width=True):
            if ARBOR_AVAILABLE:
                st.session_state.mode = "do"
                prompt = "Find free time slots for 1 hour meeting this week"
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                try:
                    result = arbor_do(prompt)
                    response = f"**Action Completed:**\n\n{result}"
                    details = {"mode": "do", "command": prompt, "timestamp": datetime.now().isoformat()}
                    st.session_state.messages.append({"role": "assistant", "content": response, "details": details})
                except Exception as e:
                    error_msg = f"**Error:** {str(e)}"
                    st.session_state.messages.append({"role": "assistant", "content": error_msg, "details": {"error": str(e)}})
                
                st.rerun()
            else:
                st.error("Arbor not available")

# Footer
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("""
<div style='text-align: center; color: #64748b; padding: 20px; 
            background: linear-gradient(135deg, rgba(79, 70, 229, 0.1) 0%, rgba(124, 58, 237, 0.1) 100%);
            border-radius: 12px; border-top: 1px solid rgba(79, 70, 229, 0.3);'>
    <strong>Arbor</strong><br>
    Powered by DeepSeek R1 & LangChain | Built for ClerkTree<br>
    <small>Advanced Analytics • Calendar Intelligence • Document Processing • Email Automation</small>
</div>
""", unsafe_allow_html=True)

# Show system status on startup
if not ARBOR_AVAILABLE:
    st.error("""
    **Setup.py not found!**
    
    Please ensure `setup.py` is in the same directory as this Streamlit app.
    The setup.py file should contain all the Arbor Enhanced functions and configurations.
    """)
elif not st.session_state.initialized:
    st.info("""
    **System Ready for Initialization**
    
    Click the "Initialize System" button in the sidebar to start using Arbor Enhanced.
    """)

# Additional quick stats at the bottom
if ARBOR_AVAILABLE and st.session_state.initialized:
    st.markdown("<br>", unsafe_allow_html=True)

    # System health check
    try:
        status = system_status()
        operational_count = sum(1 for v in status.values() if v == "✓")
        total_components = len(status)
        health_percentage = (operational_count / total_components) * 100
        
        health_color = "#10b981" if health_percentage >= 80 else "#f59e0b" if health_percentage >= 50 else "#ef4444"
        st.markdown(f"""
        <div style='background: rgba(79, 70, 229, 0.05); padding: 15px; border-radius: 12px; 
                    border: 1px solid rgba(79, 70, 229, 0.2); text-align: center;'>
            <strong style='color: #e0e7ff;'>System Health:</strong> 
            <span style='color: {health_color}; font-weight: 700; font-size: 1.2rem;'>{health_percentage:.0f}%</span>
            <span style='color: #94a3b8; font-size: 0.9rem;'> ({operational_count}/{total_components} components operational)</span>
        </div>
        """, unsafe_allow_html=True)
    except Exception as e:
        pass